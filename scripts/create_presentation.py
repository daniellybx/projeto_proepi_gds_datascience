#!/usr/bin/env python3
"""
Creates a PowerPoint presentation from notebooks 02_clustering_syndromes.ipynb
and 03_time_series_analysis.ipynb. Work in progress — preliminary results.
All content in English. Titles use "Preliminary results".
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("Install python-pptx: pip install python-pptx")

BLUE_DARK = RGBColor(0x00, 0x4B, 0x87)
GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0x66, 0x66, 0x66)


def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    left, width = Inches(0.5), Inches(9)
    tf = slide.shapes.add_textbox(left, Inches(1.2), width, Inches(1.5))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    if subtitle:
        tf2 = slide.shapes.add_textbox(left, Inches(2.9), width, Inches(1.2))
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_LIGHT
    return slide


def add_content_slide(prs, title, bullet_points, footnote=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    tf = slide.shapes.add_textbox(left, Inches(0.4), Inches(9), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    tf2 = slide.shapes.add_textbox(left, Inches(1.3), Inches(9), Inches(5.2))
    text_frame = tf2.text_frame
    text_frame.word_wrap = True
    for i, point in enumerate(bullet_points):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)
        p.level = 0
    if footnote:
        tf3 = slide.shapes.add_textbox(left, Inches(6.6), Inches(9), Inches(0.6))
        p3 = tf3.text_frame.paragraphs[0]
        p3.text = footnote
        p3.font.size = Pt(10)
        p3.font.italic = True
        p3.font.color.rgb = GRAY_LIGHT
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # —— Slide 1: Title ——
    add_title_slide(
        prs,
        "Preliminary results: From symptom reports to early alerts",
        "Participatory surveillance at UnB · Guardiões da Saúde (ProEpi) · Work in progress",
    )

    # —— Slide 2: Aims and scope ——
    add_content_slide(
        prs,
        "Preliminary results: Aims and scope",
        [
            "Participatory surveillance (Guardiões da Saúde app) generates a large stream of symptom reports from the UnB community (2022–2024).",
            "Goal: turn this stream into actionable intelligence — detect unusual patterns early and support field response.",
            "Our path: (1) find meaningful symptom clusters and link them to possible clinical entities; (2) model usual occurrence of these syndromes; (3) automate models so alerts can trigger field investigations.",
            "",
            "All results shown here are preliminary and under continuous refinement.",
        ],
        footnote="Work in progress — ProEpi Guardiões da Saúde.",
    )

    # —— Slide 3: Stage 1 — Clustering (algorithms, metrics, DBSCAN choice, 13 clusters, clinical mapping) ——
    add_content_slide(
        prs,
        "Preliminary results: Stage 1 — Clustering and 13 symptom clusters",
        [
            "Algorithms tested on cleaned symptom reports (~2M records): K-Means, GMM, DBSCAN, Agglomerative Clustering, K-Prototypes.",
            "Comparison metrics: Silhouette score, Davies-Bouldin index, Calinski-Harabasz index.",
            "",
            "DBSCAN was chosen: it showed reasonable (though not the best) evaluation metrics, but yielded the best division of symptoms — 13 clusters.",
            "",
            "We are now fitting these 13 clusters into standard clinical case definitions to identify possible diseases, conditions, aggravations, or syndromes that each cluster might represent. These are possibilities under review, not final diagnoses.",
            "",
            "Output: processed dataset with cluster labels for Stage 2.",
        ],
        footnote="Preliminary — cluster count and clinical mapping subject to revision.",
    )

    # —— Slide 4: Stage 2 — Time series (algorithms, metrics, no choice yet) ——
    add_content_slide(
        prs,
        "Preliminary results: Stage 2 — Time series modelling",
        [
            "Goal: model each (cluster-based) syndrome over time to establish a baseline pattern of occurrence.",
            "Time series models tested: Prophet, XGBoost, SARIMAX, Holt-Winters.",
            "Comparison metrics: RMSE, MAE, MAPE, MASE, directional accuracy, prediction interval coverage; diagnostic tests on residuals: Jarque-Bera, ADF.",
            "",
            "There is not yet a final model choice per cluster, given that clusters are still preliminary; model selection and validation will be finalised as the cluster–syndrome mapping is stabilised.",
            "",
            "Goal: models that can be automated and deployed so that automatic alerts are raised when observed counts exceed the expected pattern.",
        ],
        footnote="Preliminary — model selection pending stabilisation of clusters.",
    )

    # —— Slide 5: Steps that automation will trigger + closing ——
    add_content_slide(
        prs,
        "Preliminary results: What automation will trigger",
        [
            "Once models are automated and deployed in the system, the following steps will be triggered:",
            "",
            "1. Models run on incoming data (e.g. daily or weekly) to produce expected counts per syndrome.",
            "2. Observed counts are compared with expected patterns; anomalies (e.g. above a threshold or confidence interval) are detected.",
            "3. The system raises automatic alerts when an anomaly is flagged.",
            "4. Alerts trigger field investigations by the residents and interns of the UnB Health Situation Room (Sala de Situação de Saúde), closing the loop from data to action.",
            "",
            "All findings in this presentation are preliminary and are being improved as we refine methods and interpretations.",
        ],
        footnote="Work in progress — ProEpi Guardiões da Saúde · UnB.",
    )

    out_dir = Path(__file__).resolve().parent.parent
    out_path = out_dir / "apresentacao_resultados_proepi_gds.pptx"
    prs.save(str(out_path))
    print(f"Presentation saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    main()
